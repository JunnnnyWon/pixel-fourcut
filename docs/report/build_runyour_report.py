"""RunYour 후원사 제출용 14장 보고서 빌더."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[2]
GALLERY = ROOT / "docs" / "gallery"
OUT = ROOT / "docs" / "report" / "RunYour_Festival_Report_2025.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK = RGBColor(0x1F, 0x24, 0x2E)
SUB = RGBColor(0x5B, 0x63, 0x73)
ACCENT = RGBColor(0xFF, 0x6B, 0x35)
LINE = RGBColor(0xE3, 0xE6, 0xEA)
BG_CARD = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_KR = "Pretendard"
FONT_FALLBACK = "Apple SD Gothic Neo"


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT_KR):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=LINE, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def add_header(slide, title, subtitle=None, page=None, total=14):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.45),
                                 Inches(0.08), Inches(0.55))
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text(slide, Inches(0.7), Inches(0.4), Inches(10), Inches(0.55),
             title, size=24, bold=True, color=INK)
    if subtitle:
        add_text(slide, Inches(0.72), Inches(0.95), Inches(10), Inches(0.35),
                 subtitle, size=12, color=SUB)

    divider = slide.shapes.add_connector(1, Inches(0.5), Inches(1.45),
                                         Inches(12.83), Inches(1.45))
    divider.line.color.rgb = LINE
    divider.line.width = Pt(0.75)

    if page is not None:
        add_text(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.3),
                 f"{page} / {total}",
                 size=9, color=SUB, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "픽셀네컷 × RunYour  ·  청강문화산업대학교 축제 2026.05.21",
             size=9, color=SUB)


def add_bullets(slide, left, top, width, height, items, *,
                size=13, color=INK, gap=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"·  {item}"
        run.font.name = FONT_KR
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def kv_card(slide, left, top, width, height, label, value, *,
            value_size=22, label_color=SUB, value_color=INK):
    add_rect(slide, left, top, width, height, fill=BG_CARD, line=LINE)
    add_text(slide, left + Inches(0.2), top + Inches(0.18),
             width - Inches(0.4), Inches(0.3),
             label, size=10, color=label_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.5),
             width - Inches(0.4), height - Inches(0.6),
             value, size=value_size, bold=True, color=value_color)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return slide


def safe_picture(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        add_rect(slide, left, top, width or Inches(3), height or Inches(2),
                 fill=BG_CARD, line=LINE)
        return None
    return slide.shapes.add_picture(str(path), left, top,
                                    width=width, height=height)


def slide_cover(prs):
    s = new_slide(prs)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.2), SLIDE_H)
    panel.shadow.inherit = False
    panel.fill.solid()
    panel.fill.fore_color.rgb = INK
    panel.line.fill.background()

    add_text(s, Inches(0.6), Inches(0.7), Inches(4.5), Inches(0.4),
             "RunYour GPU 후원 보고서",
             size=14, color=ACCENT, bold=True)
    add_text(s, Inches(0.6), Inches(1.2), Inches(4.5), Inches(2.5),
             "픽셀네컷\n축제 운영 결과 공유",
             size=34, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(4.2), Inches(4.5), Inches(0.5),
             "청강문화산업대학교 축제 · 2026.05.21",
             size=14, color=WHITE)
    add_text(s, Inches(0.6), Inches(4.7), Inches(4.5), Inches(0.5),
             "GPU 지원사: RunYour",
             size=14, color=WHITE)
    add_text(s, Inches(0.6), Inches(6.7), Inches(4.5), Inches(0.4),
             "작성: 픽셀네컷 운영팀",
             size=10, color=RGBColor(0xB8, 0xBE, 0xC9))

    safe_picture(s, GALLERY / "banner" / "hero-photo-transform.png",
                 Inches(5.6), Inches(0.9), width=Inches(7.2))
    add_text(s, Inches(5.6), Inches(6.6), Inches(7.2), Inches(0.4),
             "AI 부스 운영 / 약 100팀 · 약 300명 참여",
             size=11, color=SUB, align=PP_ALIGN.RIGHT)


def slide_summary(prs):
    s = new_slide(prs)
    add_header(s, "한 줄 요약",
               "RunYour A40 48G 인스턴스로 축제 현장 AI 부스를 21시간 26분 무중단 운영했습니다.",
               page=2)

    top = Inches(1.9)
    h = Inches(1.5)
    w = Inches(2.95)
    gap = Inches(0.2)
    left0 = Inches(0.5)
    kv_card(s, left0 + (w + gap) * 0, top, w, h, "운영 시간", "21h 26m")
    kv_card(s, left0 + (w + gap) * 1, top, w, h, "참여 규모", "약 100팀 / 300명")
    kv_card(s, left0 + (w + gap) * 2, top, w, h, "평균 처리시간", "약 1–2분 / 컷")
    kv_card(s, left0 + (w + gap) * 3, top, w, h, "GPU 비용", "36,779원")

    add_rect(s, Inches(0.5), Inches(3.7), Inches(12.33), Inches(2.9),
             fill=BG_CARD, line=LINE)
    add_text(s, Inches(0.8), Inches(3.95), Inches(11.8), Inches(0.45),
             "결론부터 말씀드리면",
             size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.8), Inches(4.4), Inches(11.8), Inches(2.1), [
        "로컬 PC 없이도 SDXL 기반 AI 사진관을 축제 현장에서 무중단으로 굴릴 수 있었고, On-Demand 21h 26m / 36,779원으로 단발성 행사 진입 장벽이 낮았습니다.",
        "초반 약 1분이던 처리시간이 장시간 가동 후 약 2분으로 체감상 2배 늘었고, 원인은 저희도 정확히 짚지 못했습니다.",
        "운영 중 GPU 사용량·인스턴스 상태를 현장에서 빠르게 볼 수 있는 모바일/요약 뷰가 없어, 종료 시점 판단이 불안했습니다.",
        "워크플로우 셋업을 직접 할 수 있는 팀이라면, 일회성 GPU 옵션으로 합리적인 선택지였다는 게 솔직한 평가입니다.",
    ], size=13)


def slide_project(prs):
    s = new_slide(prs)
    add_header(s, "프로젝트 소개 — 픽셀네컷",
               "축제 현장에서 즉석 촬영 → AI 변환 → 2컷 인쇄까지 한 번에.",
               page=3)

    add_text(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(0.5),
             "무엇을 만들었나요?", size=15, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.5), Inches(2.2), Inches(7.5), Inches(4.5), [
        "Canon EOS 카메라로 손님을 촬영합니다.",
        "촬영본을 ComfyUI + SDXL 기반 AI 파이프라인이 자동으로 변환합니다.",
        "변환된 결과를 Canon SELPHY CP1500으로 2컷 포토카드로 출력합니다.",
        "운영 인원은 학생 1–2명이며, 손님 1팀 처리에 약 1–2분이 걸렸습니다.",
        "프론트엔드 화면에서는 촬영 → 미리보기 → 인쇄까지 한 번에 진행합니다.",
    ], size=13)

    safe_picture(s, GALLERY / "banner" / "hero-camera-booth.png",
                 Inches(8.3), Inches(1.9), width=Inches(4.6))
    add_text(s, Inches(8.3), Inches(6.55), Inches(4.6), Inches(0.3),
             "현장 부스 컨셉 일러스트", size=9, color=SUB,
             align=PP_ALIGN.CENTER)


def slide_booth(prs):
    s = new_slide(prs)
    add_header(s, "축제 현장 — 청강문화산업대학교",
               "2026.05.21, 학생 자율 부스로 운영했습니다.",
               page=4)

    safe_picture(s, GALLERY / "poster" / "pixel4cut-poster-v1.png",
                 Inches(0.5), Inches(1.8), height=Inches(5.2))
    safe_picture(s, GALLERY / "poster" / "pixel4cut-poster-v2.png",
                 Inches(4.4), Inches(1.8), height=Inches(5.2))

    add_text(s, Inches(8.3), Inches(1.8), Inches(4.5), Inches(0.5),
             "운영 개요", size=15, bold=True, color=ACCENT)
    add_bullets(s, Inches(8.3), Inches(2.3), Inches(4.7), Inches(4.5), [
        "장소: 청강문화산업대학교 축제 부스",
        "일시: 2026년 5월 21일",
        "참여: 약 100팀 / 약 300명",
        "운영 시간: 약 21시간 26분 (장시간 무중단)",
        "결과물: 2컷 포토카드 즉석 인쇄",
        "현장 사진은 학생 개인 기기에 저장되어 있어 본 보고서에는 일부만 첨부했습니다.",
    ], size=12)


def slide_pipeline(prs):
    s = new_slide(prs)
    add_header(s, "기술 파이프라인",
               "Canon EOS → 워치 폴더 → ComfyUI(SDXL) → SELPHY CP1500.",
               page=5)

    steps = [
        ("01", "촬영", "Canon EOS\n+ EOS Utility"),
        ("02", "수집", "WATCH_FOLDER\n자동 감지"),
        ("03", "AI 변환", "ComfyUI\nSDXL 워크플로우"),
        ("04", "프리뷰", "프론트엔드\n웹 UI"),
        ("05", "인쇄", "SELPHY CP1500\n2컷 포토카드"),
    ]
    top = Inches(2.1)
    h = Inches(2.0)
    w = Inches(2.36)
    gap = Inches(0.13)
    left0 = Inches(0.5)
    for i, (num, title, desc) in enumerate(steps):
        l = left0 + (w + gap) * i
        add_rect(s, l, top, w, h, fill=WHITE, line=LINE)
        add_text(s, l + Inches(0.2), top + Inches(0.18),
                 w - Inches(0.4), Inches(0.4),
                 num, size=11, color=ACCENT, bold=True)
        add_text(s, l + Inches(0.2), top + Inches(0.55),
                 w - Inches(0.4), Inches(0.5),
                 title, size=16, bold=True, color=INK)
        add_text(s, l + Inches(0.2), top + Inches(1.05),
                 w - Inches(0.4), Inches(0.9),
                 desc, size=11, color=SUB)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                l + w + Inches(0.005),
                top + Inches(0.92),
                Inches(0.12), Inches(0.18))
            arrow.rotation = 30
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
             "이 중 \"AI 변환\" 단계가 GPU를 가장 많이 씁니다. → RunYour A40 48G 인스턴스 담당.",
             size=13, color=INK, bold=True)
    add_bullets(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0), [
        "SDXL 기반 워크플로우라 VRAM 여유가 중요했습니다. A40 48G로 충분히 여유 있게 동작했습니다.",
        "촬영 → 결과 출력까지의 총 체감 시간은 평균 1–2분 수준이었습니다.",
        "로컬 PC 의존도를 낮춰서, 운영 PC가 다운돼도 변환 서버는 계속 살아 있는 구조였습니다.",
    ], size=12)


def slide_alternatives(prs):
    s = new_slide(prs)
    add_header(s, "기존 옵션과의 비교",
               "축제 한 번을 위해 어떤 GPU를 쓸 수 있었는지 솔직하게 비교했습니다.",
               page=6)

    headers = ["옵션", "초기 비용", "이번 행사 비용감", "운영 난이도", "느낀 점 (학생 입장)"]
    rows = [
        ["로컬 RTX 4090 PC",
         "약 250–300만원대",
         "장기적으로는 저렴",
         "PC 조립·발열·소음 직접 관리",
         "1회성 축제만 보면 과한 투자."],
        ["로컬 RTX 3090 PC",
         "중고 100만원대~",
         "저렴하지만 VRAM·속도 한계",
         "SDXL은 살짝 빠듯",
         "이미 갖고 있다면 OK, 새로 사긴 애매."],
        ["Vast.ai",
         "0원",
         "시간당 저렴",
         "인스턴스 품질 편차 있음",
         "스팟성/중단 가능성이 부담."],
        ["RunPod",
         "0원",
         "시간당 합리적",
         "익숙해지면 편함",
         "UI/지원은 무난, 결제 흐름이 별도."],
        ["RunYour (선택)",
         "0원 (후원)",
         "21h 26m / 36,779원",
         "On-Demand로 켜고 끄기만",
         "단발성 행사에 가장 부담이 적었음."],
    ]

    left = Inches(0.5)
    top = Inches(1.9)
    row_h = Inches(0.75)
    col_w = [Inches(2.2), Inches(2.0), Inches(2.4), Inches(2.4), Inches(3.33)]

    x = left
    for i, htxt in enumerate(headers):
        add_rect(s, x, top, col_w[i], row_h, fill=INK, line=INK)
        add_text(s, x + Inches(0.15), top + Inches(0.18),
                 col_w[i] - Inches(0.3), row_h - Inches(0.2),
                 htxt, size=11, bold=True, color=WHITE)
        x += col_w[i]

    for r, row in enumerate(rows):
        y = top + row_h * (r + 1)
        x = left
        is_us = row[0].startswith("RunYour")
        for i, cell in enumerate(row):
            add_rect(s, x, y, col_w[i], row_h,
                     fill=(ACCENT if is_us else (BG_CARD if r % 2 == 0 else WHITE)),
                     line=LINE)
            add_text(s, x + Inches(0.15), y + Inches(0.18),
                     col_w[i] - Inches(0.3), row_h - Inches(0.2),
                     cell, size=10.5,
                     color=(WHITE if is_us else INK),
                     bold=is_us and i == 0)
            x += col_w[i]

    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
             "* 비용/난이도는 학생 운영팀 기준 체감값입니다. 절대적인 우열을 가리는 표가 아닙니다.",
             size=9, color=SUB)


def slide_runyour_usage(prs):
    s = new_slide(prs)
    add_header(s, "RunYour 사용 내역",
               "이번 축제에서 실제로 결제·사용한 인스턴스입니다.",
               page=7)

    top = Inches(1.9)
    w = Inches(3.95)
    h = Inches(1.5)
    gap = Inches(0.2)
    left0 = Inches(0.5)
    kv_card(s, left0 + (w + gap) * 0, top, w, h, "인스턴스 #", "#26")
    kv_card(s, left0 + (w + gap) * 1, top, w, h, "타입", "On-Demand")
    kv_card(s, left0 + (w + gap) * 2, top, w, h, "GPU / VRAM", "A40 · 48GB")

    top2 = top + h + gap
    kv_card(s, left0 + (w + gap) * 0, top2, w, h,
            "사용 기간", "2026-05-20 20:13\n→ 2026-05-21 17:39")
    kv_card(s, left0 + (w + gap) * 1, top2, w, h, "총 시간", "21h 26m")
    kv_card(s, left0 + (w + gap) * 2, top2, w, h, "총 결제 금액", "36,779원")

    add_rect(s, Inches(0.5), Inches(5.4), Inches(12.33), Inches(1.5),
             fill=BG_CARD, line=LINE)
    add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.4),
             "왜 A40 48G였나요?",
             size=13, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.9), [
        "SDXL 워크플로우의 VRAM 여유를 우선했습니다. (모델 + LoRA + 워크플로우 동시 적재)",
        "축제 하루 동안 무중단으로 계속 켜둘 거라서, On-Demand 안정성이 중요했습니다.",
    ], size=11)


def slide_good(prs):
    s = new_slide(prs)
    add_header(s, "좋았던 점 — 솔직 후기",
               "학생 운영팀 입장에서, 다른 GPU 옵션 대비 좋았던 부분만 모았습니다.",
               page=8)

    cards = [
        ("진입 장벽이 낮음",
         "GPU 카드 구매·PC 신규 구성 없이 축제 한 번을 위한 단발성 환경을 구성할 수 있었습니다. 학생 팀 예산 범위에서 SDXL 워크플로우를 굴릴 수 있는 선택지였습니다."),
        ("48G VRAM 여유",
         "SDXL + 보조 모델 + ComfyUI 워크플로우를 한 인스턴스에 적재한 상태로 OOM 없이 운영했습니다. 워크플로우 변경 시에도 메모리 압박은 없었습니다."),
        ("On-Demand 비용 예측 가능성",
         "21h 26m 연속 가동에 36,779원으로, 사전에 예상한 범위와 큰 차이가 없었습니다. 단발성 행사 회계 처리에는 적합한 과금 구조였습니다."),
    ]
    top = Inches(1.85)
    w = Inches(4.0)
    h = Inches(4.6)
    gap = Inches(0.22)
    for i, (title, body) in enumerate(cards):
        left = Inches(0.5) + (w + gap) * i
        y = top
        add_rect(s, left, y, w, h, fill=BG_CARD, line=LINE)
        add_text(s, left + Inches(0.25), y + Inches(0.25),
                 w - Inches(0.5), Inches(0.55),
                 title, size=15, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.25), y + Inches(0.95),
                 w - Inches(0.5), h - Inches(1.1),
                 body, size=12, color=INK)


def slide_bad(prs):
    s = new_slide(prs)
    add_header(s, "불편했거나 아쉬웠던 점",
               "더 좋아지면 좋겠다는 의미로, 가감 없이 적었습니다.",
               page=9)

    cards = [
        ("장시간 가동 시 처리시간 증가",
         "초반 약 1분이던 컷당 처리시간이 시간이 지나며 약 2분까지 늘었습니다. 캐시·메모리 누수, 워크플로우 쪽 원인 등 가능성은 있으나 운영팀 차원에서 원인을 특정하지 못했습니다. 사용자 체감으로는 명확한 차이였습니다."),
        ("현장 모니터링 채널 부족",
         "축제 부스 운영 중 GPU 사용량·인스턴스 상태를 빠르게 확인할 수 있는 모바일·현장용 뷰가 없어 정량 모니터링이 불가능했습니다. 처리시간 증가의 원인 분석도 이 부재로 인해 어려웠습니다."),
        ("워크플로우 셋업은 사용자 몫",
         "ComfyUI·SDXL 환경은 직접 구성해야 했습니다. RunYour 책임 영역이 아닌 점은 이해하나, 워크플로우 셋업을 직접 할 수 있는 팀이 아니라면 진입 장벽이 됩니다. \"학생·소규모 프로젝트용 ComfyUI 템플릿\" 형태의 자산이 있다면 재사용 의향이 있습니다."),
        ("결제·세션 종료 UX",
         "처음 사용하는 입장에서 종료 시점에 \"정말 종료된 상태인지\"를 여러 번 확인하게 됐습니다. 종료 직전 요약 알림·모바일 종료 버튼 등 종료 확정 신호가 명확하면 안심감이 커질 것으로 보입니다."),
        ("비용 예측 도구 부재",
         "사전 예상 범위와 결과가 큰 차이는 없었으나, 운영 중 누적 비용을 실시간으로 확인할 수 있는 도구가 있으면 좋겠습니다. 단발성 행사 회계 처리에는 \"지금까지 얼마\"가 보이는 것이 중요합니다."),
    ]
    top = Inches(1.85)
    cols = 3
    rows = 2
    w = Inches(4.0)
    h = Inches(2.5)
    gap = Inches(0.22)
    for i, (title, body) in enumerate(cards):
        row = i // cols
        col = i % cols
        left = Inches(0.5) + (w + gap) * col
        y = top + (h + gap) * row
        add_rect(s, left, y, w, h, fill=BG_CARD, line=LINE)
        add_text(s, left + Inches(0.25), y + Inches(0.2),
                 w - Inches(0.5), Inches(0.45),
                 title, size=14, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.25), y + Inches(0.75),
                 w - Inches(0.5), h - Inches(0.9),
                 body, size=10, color=INK)


def slide_metrics(prs):
    s = new_slide(prs)
    add_header(s, "운영 지표",
               "현장에서 정확히 카운트한 값과, 추정값을 구분해서 적었습니다.",
               page=10)

    top = Inches(1.9)
    w = Inches(2.95)
    h = Inches(1.5)
    gap = Inches(0.2)
    left0 = Inches(0.5)
    kv_card(s, left0 + (w + gap) * 0, top, w, h, "참여 팀수", "약 100팀")
    kv_card(s, left0 + (w + gap) * 1, top, w, h, "총 참여 인원", "약 300명")
    kv_card(s, left0 + (w + gap) * 2, top, w, h, "처리시간 (초반)", "약 1분 / 컷")
    kv_card(s, left0 + (w + gap) * 3, top, w, h, "처리시간 (후반)", "약 2분 / 컷")

    add_rect(s, Inches(0.5), Inches(3.7), Inches(12.33), Inches(3.2),
             fill=BG_CARD, line=LINE)
    add_text(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.45),
             "정확히 카운트하지 못한 지표 (개선 과제)",
             size=14, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.4), [
        "실제 인쇄 컷 수: 현장에서 카운트를 따로 하지 못했습니다. 100팀 × 2컷 가정 시 약 200컷이 상한 추정치이며, 실제 수치는 확인 불가입니다.",
        "동시 접속·대기 팀 수: 줄 길이 기반 체감만 존재해 정량 수치로 기록하지 못했습니다.",
        "GPU 사용률·메모리 점유: 운영 중 실시간 로그를 수집하지 못했고, 처리시간 증가의 원인 분석도 이로 인해 막혔습니다.",
        "세션별 처리시간 분포: 평균 외 분산·최대값을 별도 기록하지 못해 \"체감 2배\"라는 정성 진술만 남았습니다.",
        "다음 운영 시 인쇄 카운터, 대기열 카운터, GPU·메모리 로깅, 세션별 시간 측정을 기본 인프라로 포함할 계획입니다.",
    ], size=12)


def slide_proud(prs):
    s = new_slide(prs)
    add_header(s, "자랑하고 싶은 부분",
               "운영팀 입장에서 \"이건 잘 만든 것 같아요\" 항목입니다.",
               page=11)

    add_bullets(s, Inches(0.5), Inches(1.9), Inches(8.0), Inches(5.0), [
        "카메라 → AI → 인쇄까지 한 사람이 1–2분 안에 끝낼 수 있는 워크플로우를 학생 팀이 직접 만들었습니다.",
        "외부 GPU에 의존하면서도, 현장에서는 \"인터넷 약간 느린 사진관\" 정도의 체감으로 동작했습니다.",
        "포스터·현수막·프롬프트 보드까지, 부스 비주얼을 자체 디자인 자산으로 통일했습니다.",
        "축제 하루 동안 약 100팀 / 300명을 별다른 큰 사고 없이 응대했습니다.",
        "운영 종료 후 모든 결과물을 정리해 오픈소스 레포에 공개했습니다. (얼굴 사진은 비공개)",
    ], size=13)

    safe_picture(s, GALLERY / "banner" / "hero-photo-transform.png",
                 Inches(8.6), Inches(1.9), width=Inches(4.4))
    add_text(s, Inches(8.6), Inches(6.4), Inches(4.4), Inches(0.4),
             "AI 변환 컨셉 일러스트", size=9, color=SUB,
             align=PP_ALIGN.CENTER)


def slide_gallery(prs):
    s = new_slide(prs)
    add_header(s, "현장 자산 갤러리",
               "축제에서 사용한 현수막·포스터 자산입니다. (손님 얼굴 사진은 비공개)",
               page=12)

    items = [
        ("hero-camera-booth.png",      "부스 메인"),
        ("hero-photo-transform.png",   "AI 변환 컨셉"),
        ("title-plaque.png",           "타이틀 명패"),
        ("prompts-plaque.png",         "프롬프트 보드"),
        ("bg-sky-base.png",            "현수막 배경 (하늘)"),
        ("bg-ground-strip.png",        "현수막 배경 (지면)"),
        ("cloud-01.png",               "구름 데코"),
        ("sparkle-pack.png",           "스파클 데코"),
    ]
    left0 = Inches(0.5)
    top0 = Inches(1.9)
    cell_w = Inches(3.0)
    cell_h = Inches(2.45)
    gap_x = Inches(0.13)
    gap_y = Inches(0.18)
    for i, (fname, caption) in enumerate(items):
        col = i % 4
        row = i // 4
        x = left0 + (cell_w + gap_x) * col
        y = top0 + (cell_h + gap_y) * row
        add_rect(s, x, y, cell_w, cell_h, fill=BG_CARD, line=LINE)
        safe_picture(s, GALLERY / "banner" / fname,
                     x + Inches(0.12), y + Inches(0.12),
                     width=cell_w - Inches(0.24))
        add_text(s, x + Inches(0.12),
                 y + cell_h - Inches(0.42),
                 cell_w - Inches(0.24), Inches(0.3),
                 caption, size=10, color=SUB,
                 align=PP_ALIGN.CENTER)


def slide_future(prs):
    s = new_slide(prs)
    add_header(s, "다음 학기 / 다음 행사 계획",
               "지금 상태에서 한 단계만 더 가면 좋겠다 싶은 항목입니다.",
               page=13)

    add_bullets(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(5.0), [
        "장시간 가동 시 컷당 처리시간이 늘어나는 원인을 찾아서 캐시/메모리 관리 로직을 정리할 계획입니다.",
        "현장 운영용 대시보드 (GPU 사용률 / 대기열 / 인쇄 카운트)를 모바일 친화적으로 만들고 싶습니다.",
        "워크플로우를 \"학생용 ComfyUI 템플릿\"으로 정리해서, 다음 학기 후배들이 그대로 쓸 수 있게 만들 계획입니다.",
        "다음 행사에서도 RunYour를 후보 1순위로 두고, 동일 인스턴스로 한 번 더 검증할 생각입니다.",
        "오픈소스 레포는 계속 유지·정리하고, 필요 시 RunYour 케이스 스터디로 함께 노출되어도 좋다는 입장입니다.",
    ], size=13)


def slide_thanks(prs):
    s = new_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(3.0),
                             Inches(0.1), Inches(1.5))
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text(s, Inches(0.85), Inches(2.8), Inches(11), Inches(0.5),
             "감사합니다.",
             size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.85), Inches(3.9), Inches(11), Inches(0.5),
             "RunYour의 GPU 지원 덕분에 학생 프로젝트가 축제 현장까지 무사히 갈 수 있었습니다.",
             size=16, color=WHITE)
    add_text(s, Inches(0.85), Inches(4.5), Inches(11), Inches(0.5),
             "다음 운영에서도, 다음 학생 팀에게도 좋은 옵션이 되어주시길 바랍니다.",
             size=16, color=WHITE)

    add_text(s, Inches(0.85), Inches(6.6), Inches(11), Inches(0.4),
             "픽셀네컷 운영팀  ·  청강문화산업대학교 축제 2026.05.21  ·  14 / 14",
             size=10, color=RGBColor(0xB8, 0xBE, 0xC9))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_summary(prs)
    slide_project(prs)
    slide_booth(prs)
    slide_pipeline(prs)
    slide_alternatives(prs)
    slide_runyour_usage(prs)
    slide_good(prs)
    slide_bad(prs)
    slide_metrics(prs)
    slide_proud(prs)
    slide_gallery(prs)
    slide_future(prs)
    slide_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"OK: {OUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
